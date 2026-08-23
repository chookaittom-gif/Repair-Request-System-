from pathlib import Path
import textwrap
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


BASE = Path(r"D:\Webapp\ระบบแจ้งซ่อม")
IMG_DIR = BASE / "Pic คู่มือ"
LOGO = BASE / "fix" / "คู่มือ" / "SDULP-FIX.png"
OUT_DIR = BASE / "Presentation"
PPTX_OUT = OUT_DIR / "ระบบแจ้งซ่อมออนไลน์_2.0.pptx"
SCRIPT_OUT = OUT_DIR / "ระบบแจ้งซ่อมออนไลน์_2.0_คำบรรยาย.md"

FONT = "TH Sarabun"
W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(17, 34, 64)
MUTED = RGBColor(87, 99, 120)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
RED = RGBColor(220, 38, 38)
GREEN = RGBColor(16, 185, 129)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(226, 232, 240)


def img(name: str) -> Path:
    return IMG_DIR / name


def add_textbox(slide, x, y, w, h, text, size=24, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bg(slide, accent=BLUE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, Inches(0.62))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT
    band.line.fill.background()
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.62), W, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    slide.shapes.add_picture(str(LOGO), W - Inches(1.04), Inches(0.12), width=Inches(0.54))


def add_title(slide, title, subtitle=None, accent=BLUE):
    add_bg(slide, accent)
    add_textbox(slide, Inches(0.65), Inches(0.95), Inches(8.6), Inches(0.45), title, 30, NAVY, True)
    if subtitle:
        add_textbox(slide, Inches(0.67), Inches(1.38), Inches(8.8), Inches(0.3), subtitle, 18, MUTED)


def add_page(slide, n):
    add_textbox(slide, Inches(6.0), Inches(6.93), Inches(1.3), Inches(0.35), str(n), 26, NAVY, True, PP_ALIGN.CENTER)


def fit_picture(slide, path, x, y, w, h, radius=True):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, px - Inches(0.03), py - Inches(0.03), pw + Inches(0.06), ph + Inches(0.06))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(241, 245, 249)
    shadow.line.color.rgb = BORDER
    pic = slide.shapes.add_picture(str(path), px, py, width=pw, height=ph)
    return pic, (px, py, pw, ph)


def add_callout(slide, label, target, box, color=BLUE):
    x, y, w, h = box
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.4)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = NAVY
    x1 = x + (w / 2)
    y1 = y + h
    x2, y2 = target
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return shape


def add_section_label(slide, text, x, y, color=TEAL):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(2.7), Inches(0.34))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf = pill.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    return pill


def image_grid(slide, names, x, y, w, h, cols=2, gap=Inches(0.12)):
    cells = []
    rows = (len(names) + cols - 1) // cols
    cw = (w - gap * (cols - 1)) / cols
    ch = (h - gap * (rows - 1)) / rows
    for i, name in enumerate(names):
        col = i % cols
        row = i // cols
        cells.append(fit_picture(slide, img(name), x + col * (cw + gap), y + row * (ch + gap), cw, ch)[1])
    return cells


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(246, 248, 252)
    side_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(4.0), H)
    side_band.fill.solid()
    side_band.fill.fore_color.rgb = RGBColor(232, 240, 255)
    side_band.line.fill.background()
    slide.shapes.add_picture(str(LOGO), W - Inches(1.2), Inches(0.25), width=Inches(0.72))
    add_textbox(slide, Inches(0.9), Inches(1.8), Inches(9.4), Inches(0.8), "ระบบแจ้งซ่อมออนไลน์ 2.0", 46, NAVY, True)
    add_textbox(slide, Inches(0.92), Inches(2.72), Inches(7.8), Inches(0.4), "สรุปภาพรวมการใช้งานสำหรับบุคคลทั่วไป เจ้าหน้าที่ และผู้บริหาร", 24, MUTED)
    add_textbox(slide, Inches(0.92), Inches(4.85), Inches(8.6), Inches(0.8), "พัฒนาโดย นายชูเกียรติ กุศลสถิตย์\nมหาวิทยาลัยสวนดุสิต ศูนย์การศึกษาลำปาง", 24, NAVY, False)
    fit_picture(slide, img("1.jpg"), Inches(7.1), Inches(1.55), Inches(5.35), Inches(3.25))


def add_slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "แจ้งซ่อมง่ายๆ สำหรับบุคคลทั่วไป", "กรอกข้อมูล แนบรูป และส่งถึงทีมช่างในขั้นตอนเดียว")
    cells = image_grid(slide, ["1.jpg", "4.jpg"], Inches(0.7), Inches(1.8), Inches(7.15), Inches(4.85), cols=2)
    add_callout(slide, "สร้างเลขที่/วันที่ให้อัตโนมัติ", (cells[0][0] + cells[0][2] * 0.47, cells[0][1] + cells[0][3] * 0.55), (Inches(8.35), Inches(1.95), Inches(3.55), Inches(0.55)))
    add_callout(slide, "แนบรูปปัญหาให้ช่างดูได้ทันที", (cells[1][0] + cells[1][2] * 0.55, cells[1][1] + cells[1][3] * 0.63), (Inches(8.35), Inches(3.0), Inches(3.55), Inches(0.55)), TEAL)
    add_callout(slide, "คลิกเดียวส่งข้อมูลถึงทีมช่าง", (cells[1][0] + cells[1][2] * 0.55, cells[1][1] + cells[1][3] * 0.88), (Inches(8.35), Inches(4.05), Inches(3.55), Inches(0.55)), GREEN)
    add_page(slide, 2)


def add_simple_slide(prs, page, title, subtitle, names, callouts, accent=BLUE, cols=2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle, accent)
    cells = image_grid(slide, names, Inches(0.7), Inches(1.75), Inches(7.35), Inches(4.9), cols=cols)
    colors = [BLUE, TEAL, GREEN, RED]
    for i, (label, cx, cy, bx, by) in enumerate(callouts):
        cell = cells[min(i, len(cells) - 1)]
        target = (cell[0] + cell[2] * cx, cell[1] + cell[3] * cy)
        add_callout(slide, label, target, (Inches(bx), Inches(by), Inches(3.65), Inches(0.58)), colors[i % len(colors)])
    add_page(slide, page)
    return slide


def final_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, TEAL)
    add_textbox(slide, Inches(0.75), Inches(1.15), Inches(7.2), Inches(0.7), "ข้อมูลการติดต่อ (Support)", 40, NAVY, True)
    add_textbox(slide, Inches(0.78), Inches(2.15), Inches(9.2), Inches(0.42), "เข้าใช้งานระบบได้ที่: https://108links.com/sdulpfix", 24, BLUE, True)
    add_textbox(slide, Inches(0.78), Inches(3.05), Inches(8.8), Inches(1.45), "พบปัญหาการใช้งาน ติดต่อผู้ดูแลระบบ:\nนายชูเกียรติ กุศลสถิตย์\nอีเมล: chookait_kus@dusit.ac.th", 26, NAVY)
    add_textbox(slide, Inches(0.78), Inches(6.55), Inches(11.8), Inches(0.35), "มหาวิทยาลัยสวนดุสิต ศูนย์การศึกษาลำปาง | ระบบแจ้งซ่อมออนไลน์ 2.0 © 2026", 20, MUTED, False, PP_ALIGN.CENTER)


NARRATION = [
    ("Slide 1", "สวัสดีครับ/ค่ะ วันนี้ขอนำเสนอระบบแจ้งซ่อมออนไลน์ 2.0 ของมหาวิทยาลัยสวนดุสิต ศูนย์การศึกษาลำปาง ระบบนี้ช่วยให้การแจ้งซ่อม ติดตามงาน จัดการอะไหล่ และออกรายงานทำได้รวดเร็วขึ้นในที่เดียว"),
    ("Slide 2", "สำหรับบุคคลทั่วไป การแจ้งซ่อมเริ่มจากระบบสร้างเลขที่และวันที่ให้อัตโนมัติ ผู้แจ้งกรอกข้อมูลปัญหา แนบรูปภาพประกอบ แล้วกดส่งแจ้งซ่อมเพียงครั้งเดียว ข้อมูลจะถูกส่งถึงทีมช่างทันที"),
    ("Slide 3", "หลังส่งเรื่องแล้ว ผู้ใช้งานสามารถเปิดดูรายการแจ้งซ่อมและตรวจสอบป้ายสถานะของงานตนเองได้ตลอดเวลา จึงเห็นความคืบหน้าแบบเรียลไทม์โดยไม่ต้องโทรสอบถามซ้ำ"),
    ("Slide 4", "ส่วนเจ้าหน้าที่เข้าสู่ระบบผ่านหน้าล็อกอิน เพื่อรักษาความปลอดภัยของข้อมูล ผู้เกี่ยวข้องจึงเข้าถึงการจัดการงานได้ตามสิทธิ์ที่เหมาะสม"),
    ("Slide 5", "หน้าแดชบอร์ดช่วยสรุปภาพรวมงานทั้งหมด ทั้งจำนวนงานแต่ละสถานะ กราฟประเภทงาน และภาระงานของช่างแต่ละคน ทำให้ผู้ดูแลระบบเห็นสถานการณ์ได้รวดเร็ว"),
    ("Slide 6", "เมื่อมีงานเข้ามา เจ้าหน้าที่หรือช่างสามารถดูรายละเอียด เปลี่ยนสถานะ และมอบหมายผู้รับผิดชอบได้ทันที ระบบจึงช่วยให้การจ่ายงานและติดตามงานเป็นระเบียบมากขึ้น"),
    ("Slide 7", "ระบบมีหน้าต่างยืนยันก่อนลบข้อมูล โดยใช้ปุ่มยืนยันสีแดงและข้อความเตือนชัดเจน ช่วยลดความผิดพลาดและป้องกันข้อมูลสำคัญสูญหาย"),
    ("Slide 8", "เมื่อมีงานใหม่หรือมีการเปลี่ยนสถานะ ระบบส่งแจ้งเตือนผ่าน Telegram ให้ทีมช่างทันที พร้อมข้อมูลสำคัญและรูปภาพประกอบ ช่วยให้ทีมรับทราบงานจากมือถือได้เร็วขึ้น"),
    ("Slide 9", "ระบบสต็อกช่วยดูรายการอะไหล่คงเหลือและทำรายการรับเข้า เบิกออก หรือใช้ในงานซ่อมได้ เมื่อบันทึกรายการ ระบบจะปรับยอดคงเหลือให้อัตโนมัติ"),
    ("Slide 10", "นอกจากยอดคงเหลือ ระบบยังแสดงประวัติการรับและเบิกจ่ายย้อนหลัง พร้อมเตือนเมื่ออะไหล่ใกล้หมด ทำให้ตรวจสอบได้โปร่งใสและเตรียมจัดซื้อได้ทันเวลา"),
    ("Slide 11", "ระบบสามารถออกรายงาน PDF อัตโนมัติ ทั้งรายงานแจ้งซ่อม รายละเอียดสถานะ และรูปภาพประกอบ รายงานพร้อมนำไปใช้งานหรือส่งต่อได้ทันที"),
    ("Slide 12", "สำหรับรายงานประจำเดือน ระบบส่งไฟล์ PDF เข้า Telegram อัตโนมัติทุกวันที่ 25 ของเดือน ผู้บริหารจึงเปิดดูรายงานจากมือถือได้สะดวก"),
    ("Slide 13", "หากต้องการเข้าใช้งานระบบ สามารถเปิดผ่านลิงก์ https://108links.com/sdulpfix และหากพบปัญหาในการใช้งาน สามารถติดต่อผู้ดูแลระบบตามข้อมูลที่แสดงบนหน้าจอ"),
]


def write_script():
    lines = ["# คำบรรยายนำเสนอ: ระบบแจ้งซ่อมออนไลน์ 2.0", ""]
    for title, text in NARRATION:
        lines.extend([f"## {title}", text, ""])
    SCRIPT_OUT.write_text("\n".join(lines), encoding="utf-8")


def main():
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    required = [LOGO] + [img(f"{i}.jpg") for i in range(1, 25)] + [img("10-1.jpg")]
    missing = [str(p) for p in required if not p.exists()]
    if missing:
        raise FileNotFoundError("Missing required files: " + ", ".join(missing))

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    cover(prs)
    add_slide_two(prs)
    add_simple_slide(prs, 3, "ติดตามสถานะเรียลไทม์", "ตรวจสอบความคืบหน้าของงานได้ด้วยตนเอง", ["5.jpg"], [("เช็คความคืบหน้างานของตนเองได้ตลอดเวลา", 0.86, 0.30, 8.35, 2.4)], TEAL, 1)
    add_simple_slide(prs, 4, "เข้าสู่ระบบสำหรับเจ้าหน้าที่", "จัดการข้อมูลผ่านระบบที่มีการยืนยันตัวตน", ["2.jpg", "3.jpg"], [("ระบบรักษาความปลอดภัยในการจัดการข้อมูล", 0.52, 0.72, 8.35, 2.25), ("รีเซ็ตรหัสผ่านผ่านอีเมลที่ลงทะเบียน", 0.55, 0.47, 8.35, 3.35)], BLUE, 2)
    add_simple_slide(prs, 5, "แดชบอร์ดสรุปผล", "ดูภาพรวมงานและภาระงานของทีมได้ในหน้าเดียว", ["6.jpg", "7.jpg"], [("สรุปภาพรวมงานทั้งหมดแยกตามสถานะ", 0.45, 0.20, 8.35, 2.05), ("สรุปประเภทงานและภาระงานช่างแต่ละคน", 0.55, 0.37, 8.35, 3.2)], TEAL, 2)
    add_simple_slide(prs, 6, "จัดการงานและอัปเดตสถานะ", "อัปเดตความคืบหน้าและมอบหมายงานได้ทันที", ["10.jpg", "12.jpg", "13.jpg", "14.jpg"], [("ช่างอัปเดตความคืบหน้าได้ทันที", 0.53, 0.28, 8.35, 1.9), ("เปิดดูรายละเอียดก่อนดำเนินการ", 0.42, 0.38, 8.35, 2.8), ("ดูรูปประกอบงานซ่อมได้ชัดเจน", 0.54, 0.52, 8.35, 3.7), ("จ่ายงานให้ช่างได้ตรงจุด", 0.55, 0.50, 8.35, 4.6)], BLUE, 2)
    add_simple_slide(prs, 7, "ระบบป้องกันความผิดพลาด", "ยืนยันก่อนลบทุกครั้ง ลดความเสี่ยงข้อมูลสูญหาย", ["10-1.jpg"], [("แจ้งเตือนก่อนลบเสมอ ป้องกันข้อมูลหาย", 0.68, 0.66, 8.35, 2.7)], RED, 1)
    add_simple_slide(prs, 8, "แจ้งเตือนผ่าน Telegram", "งานใหม่และการเปลี่ยนสถานะส่งถึงมือถือทีมช่างทันที", ["14.jpg"], [("แจ้งงานใหม่และสถานะเข้ามือถือทีมช่างทันที", 0.18, 0.35, 8.35, 2.45)], TEAL, 1)
    add_simple_slide(prs, 9, "จัดการสต็อกและเบิกจ่าย", "ควบคุมอะไหล่เข้า-ออกให้แม่นยำและตรวจสอบง่าย", ["15.jpg", "18.jpg", "19.jpg", "21.jpg"], [("เช็คอะไหล่เข้า-ออกแม่นยำ", 0.48, 0.78, 8.35, 1.85), ("ยืนยันก่อนลบข้อมูลสต็อก", 0.52, 0.54, 8.35, 2.75), ("ทำรายการรับหรือเบิกได้ทันที", 0.40, 0.37, 8.35, 3.65), ("ตัดสต็อกอัตโนมัติเมื่อใช้ซ่อม", 0.64, 0.36, 8.35, 4.55)], BLUE, 2)
    add_simple_slide(prs, 10, "ประวัติและแจ้งเตือนอะไหล่", "ตรวจสอบย้อนหลังและรู้ก่อนเมื่ออะไหล่ใกล้หมด", ["16.jpg", "17.jpg", "20.jpg", "22.jpg"], [("บันทึกข้อมูลอะไหล่และจุดเตือน", 0.47, 0.50, 8.35, 1.8), ("แก้ไขจุดเตือนได้ตามจริง", 0.50, 0.58, 8.35, 2.7), ("เตือนล่วงหน้าเมื่ออะไหล่ใกล้หมด", 0.50, 0.70, 8.35, 3.6), ("ตรวจสอบประวัติเบิกจ่ายย้อนหลังได้", 0.77, 0.26, 8.35, 4.5)], RED, 2)
    add_simple_slide(prs, 11, "รายงานสรุปอัตโนมัติ", "ออกรายงาน PDF พร้อมข้อมูลและรูปภาพประกอบ", ["8.jpg", "9.jpg", "23.jpg"], [("ข้อมูลใน PDF พร้อมนำไปใช้ได้ทันที", 0.56, 0.39, 8.35, 2.0), ("รายงานรวมสถานะและรูปภาพประกอบ", 0.45, 0.69, 8.35, 3.0), ("ประวัติสต็อกตรวจสอบย้อนหลังได้", 0.78, 0.28, 8.35, 4.0)], TEAL, 2)
    add_simple_slide(prs, 12, "ส่งรายงานอัตโนมัติประจำเดือน", "ระบบส่งรายงาน PDF เข้า Telegram ทุกวันที่ 25", ["24.jpg"], [("ระบบส่งรายงานทุกวันที่ 25 ของเดือนอัตโนมัติ", 0.14, 0.43, 8.35, 2.35), ("ผู้บริหารกดเปิดดูจากมือถือได้เลย", 0.22, 0.22, 8.35, 3.45)], BLUE, 1)
    final_slide(prs)
    prs.save(PPTX_OUT)
    write_script()
    print(PPTX_OUT)
    print(SCRIPT_OUT)


if __name__ == "__main__":
    main()
